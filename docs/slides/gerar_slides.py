# -*- coding: utf-8 -*-
"""
Gerador de slides — RabbitMQ em Sistemas Distribuídos
Tema claro (estilo EXEMPLO), capa e conclusão escuras, Century Schoolbook para títulos.
25 slides | python gerar_slides.py  ->  RabbitMQ_Apresentacao.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os

# ════════════════════════════════════════════════════════════════════════════
# PALETA  (extraída do EXEMPLO-RabbitMQ_SistemasDistribuido.pptx)
# ════════════════════════════════════════════════════════════════════════════
def C(hexstr):
    return RGBColor.from_string(hexstr)

BG    = C("F4F5F7")   # fundo dos slides de conteúdo
CV    = C("1C1B22")   # fundo da capa / conclusão
WH    = C("FFFFFF")   # cards / superfícies
TXT   = C("1C1B22")   # texto principal / títulos
SUBT  = C("51606E")   # subtítulo do slide / corpo sob heading colorido
MUT   = C("6C7884")   # corpo de card / código / numeração de página
OR    = C("E8541E")   # laranja (acento principal)
TL    = C("1F7A8C")   # teal
AM    = C("B9770E")   # âmbar
GN    = C("2F8F5B")   # verde
RD    = C("C0392B")   # vermelho (não suportado)
BD    = C("D7DBE0")   # borda de card
SHD   = "8A8F98"      # cor da sombra (hex puro, usado em XML)
CVSUB = C("C9CDD4")   # subtítulo da capa
CVMUT = C("9AA0AA")   # texto secundário da capa
CVDIV = C("3A3A45")   # linha divisória da capa
ROWA  = WH
ROWB  = C("EDEEF0")
CODEBG = C("EEF2F2")
WARNBG = C("FBF1E2")

SERIF = "Century Schoolbook"
SANS  = "Calibri"
MONO  = "Courier New"

L, RM = 0.70, 12.63          # margens de conteúdo
CW    = RM - L               # largura útil de conteúdo
CONTENT_Y = 2.30             # topo padrão da área de conteúdo
BOTTOM    = 6.95             # limite inferior seguro (antes do número de página)

ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)), "assets")
OUTPUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "RabbitMQ_Apresentacao.pptx")
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def new_slide():
    return prs.slides.add_slide(BLANK)


# ════════════════════════════════════════════════════════════════════════════
# PRIMITIVOS
# ════════════════════════════════════════════════════════════════════════════
def _zero_margins(tf, anchor=None, wrap=True):
    bp = tf._txBody.bodyPr
    bp.set('lIns', '0'); bp.set('tIns', '0'); bp.set('rIns', '0'); bp.set('bIns', '0')
    bp.set('wrap', 'square' if wrap else 'none')
    if anchor:
        bp.set('anchor', anchor)
    # Box dimensions are always computed explicitly in this script, so disable
    # any autosize-to-text behaviour (default python-pptx textboxes ship with
    # <a:spAutoFit/>, which combined with wrap="none" causes PowerPoint/LibreOffice
    # to shrink the shape to one line and re-anchor it -- visually it looks
    # "centered" even though paragraph alignment is set to LEFT).
    for tag in ('a:spAutoFit', 'a:normAutofit', 'a:noAutofit'):
        el = bp.find(qn(tag))
        if el is not None:
            bp.remove(el)
    etree.SubElement(bp, qn('a:noAutofit'))


def shadow(shape, color=SHD, alpha="18000", blur="101600", dist="25400", direction="5400000"):
    spPr = shape._element.spPr
    for el in spPr.findall(qn('a:effectLst')):
        spPr.remove(el)
    eff = etree.SubElement(spPr, qn('a:effectLst'))
    sh = etree.SubElement(eff, qn('a:outerShdw'))
    sh.set('blurRad', blur); sh.set('dist', dist); sh.set('dir', direction)
    sh.set('algn', 'bl'); sh.set('rotWithShape', '0')
    c = etree.SubElement(sh, qn('a:srgbClr')); c.set('val', color)
    a = etree.SubElement(c, qn('a:alpha')); a.set('val', alpha)


def R_(sl, l, t, w, h, fill=None, line=None, lw=1.0, sh=False):
    """Retângulo simples."""
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.shadow.inherit = False
    if fill is not None:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line is not None:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    if sh:
        shadow(s)
    return s


def OVAL(sl, l, t, w, h, fill=None, line=None, lw=1.0):
    s = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(w), Inches(h))
    s.shadow.inherit = False
    if fill is not None:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line is not None:
        s.line.color.rgb = line; s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s


def LINE(sl, x1, y1, x2, y2, color=SUBT, w=1.5, arrow_end=True, dash=None):
    """Linha/seta real via connector (sempre fica centrada, nunca usa glifo de texto)."""
    conn = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    conn.line.color.rgb = color
    conn.line.width = Pt(w)
    ln = conn.line._get_or_add_ln()
    if arrow_end:
        tail = etree.SubElement(ln, qn('a:tailEnd'))
        tail.set('type', 'triangle')
    if dash:
        d = etree.SubElement(ln, qn('a:prstDash')); d.set('val', dash)
    return conn


def T(sl, text, l, t, w, h, sz=13, bold=False, italic=False, color=TXT,
      align=PP_ALIGN.LEFT, font=SANS, anchor=None, wrap=True, spc=None):
    tx = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame
    _zero_margins(tf, anchor=anchor, wrap=wrap)
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font
    if spc is not None:
        r._r.get_or_add_rPr().set('spc', str(spc))
    return tx


def PARA(sl, lines, l, t, w, h, sz=13, color=TXT, font=SANS, bold=False,
         italic=False, lh=1.2, align=PP_ALIGN.LEFT, anchor=None):
    """Várias linhas/parágrafos numa única textbox (com word-wrap natural)."""
    if isinstance(lines, str):
        lines = [lines]
    tx = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame; tf.word_wrap = True
    _zero_margins(tf, anchor=anchor)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = lh
        r = p.add_run(); r.text = ln if ln else " "
        r.font.size = Pt(sz); r.font.color.rgb = color; r.font.name = font
        r.font.bold = bold; r.font.italic = italic
    return tx


def T2(sl, runs, l, t, w, h, align=PP_ALIGN.LEFT, lh=1.2, anchor=None):
    """Texto com múltiplos estilos numa linha: runs = [{'text','bold','italic','color','sz','font'}]."""
    tx = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame; tf.word_wrap = True
    _zero_margins(tf, anchor=anchor)
    p = tf.paragraphs[0]; p.alignment = align; p.line_spacing = lh
    for spec in runs:
        r = p.add_run(); r.text = spec.get('text', '')
        r.font.size = Pt(spec.get('sz', 13)); r.font.bold = spec.get('bold', False)
        r.font.italic = spec.get('italic', False)
        r.font.color.rgb = spec.get('color', TXT)
        r.font.name = spec.get('font', SANS)
    return tx


def bgfill(sl, color):
    R_(sl, -0.05, -0.05, 13.43, 7.6, fill=color)


def pagenum(sl, n):
    T(sl, str(n), 12.23, 7.05, 0.70, 0.30, sz=10, color=MUT, align=PP_ALIGN.RIGHT, anchor='ctr')


def badge(sl, text, color=OR, y=0.65):
    R_(sl, L, y, 0.13, 0.13, fill=color)
    T(sl, text.upper(), L + 0.24, y - 0.085, 9.5, 0.30, sz=11.5, bold=True,
      color=color, font=SANS, anchor='ctr', spc=220)


def head(sl, sec, title, sub=None, badge_color=OR, title_sz=33, light=True):
    bgfill(sl, BG if light else CV)
    if sec:
        badge(sl, sec, color=badge_color)
    T(sl, title, L - 0.02, 0.88, CW + 0.05, 0.82, sz=title_sz, bold=True,
      color=(TXT if light else WH), font=SERIF, anchor='t', wrap=True)
    if sub:
        T(sl, sub, L, 1.68, CW, 0.50, sz=14, italic=True,
          color=(SUBT if light else CVSUB), font=SANS, anchor='t')


def card(sl, l, t, w, h, border=BD, fill=WH, sh=True, lw=1.0):
    return R_(sl, l, t, w, h, fill=fill, line=border, lw=lw, sh=sh)


def bullets(sl, items, l, t, w, sz=13, dot_color=OR, text_color=SUBT, dy=0.50, font=SANS):
    y = t
    for it in items:
        R_(sl, l, y + 0.075, 0.09, 0.09, fill=dot_color)
        T(sl, it, l + 0.26, y, w - 0.26, dy - 0.04, sz=sz, color=text_color, font=font, anchor='t')
        y += dy
    return y


def code(sl, l, t, w, h, lines, sz=11, color=TL, fill=CODEBG):
    R_(sl, l, t, w, h, fill=fill)
    PARA(sl, lines.strip('\n').split('\n'), l + 0.18, t + 0.12, w - 0.36, h - 0.24,
         sz=sz, color=color, font=MONO, lh=1.3)


def divider_h(sl, l, t, w, color=BD):
    LINE(sl, l, t, l + w, t, color=color, w=0.75, arrow_end=False)


def divider_v(sl, x, t, h, color=BD):
    LINE(sl, x, t, x, t + h, color=color, w=0.75, arrow_end=False)


def picture_fit(sl, path, l, t, max_w, max_h, align='center'):
    """Insere imagem mantendo proporção, encaixada em max_w x max_h, centrada."""
    from PIL import Image
    iw, ih = Image.open(path).size
    ratio = iw / ih
    w, h = max_w, max_w / ratio
    if h > max_h:
        h = max_h
        w = max_h * ratio
    if align == 'center':
        x = l + (max_w - w) / 2
    elif align == 'left':
        x = l
    else:
        x = l + (max_w - w)
    y = t + (max_h - h) / 2
    sl.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    return x, y, w, h


def placeholder(sl, l, t, w, h, tipo, label, instrucao=None):
    R_(sl, l, t, w, h, fill=C("FCEEE6"), line=OR, lw=1.2)
    pad = 0.18
    T(sl, "[ %s \u00b7 SUBSTITUIR ]" % tipo.upper(), l + pad, t + 0.12, w - 2 * pad, 0.26,
      sz=9.5, bold=True, color=OR, spc=60)
    if h > 0.55:
        T(sl, label, l + pad, t + 0.40, w - 2 * pad, 0.55, sz=12, bold=True, color=TXT, wrap=True)
    if h > 1.0 and instrucao:
        PARA(sl, instrucao if isinstance(instrucao, list) else [instrucao],
             l + pad, t + 0.90, w - 2 * pad, h - 1.0, sz=10, color=MUT, italic=True, lh=1.25)


def text_block(sl, title, desc, l, t, w, title_color=TXT, title_sz=14, body_sz=12.5):
    """Bloco título+corpo. Regra do design: título colorido -> corpo escuro; título escuro -> corpo cinza."""
    body_color = TXT if title_color != TXT else MUT
    T(sl, title, l, t, w, 0.32, sz=title_sz, bold=True, color=title_color, anchor='t')
    PARA(sl, desc, l, t + 0.36, w, 0.95, sz=body_sz, color=body_color, lh=1.22)


def dotcard(sl, l, t, w, h, dot_color, title, desc, title_sz=14):
    card(sl, l, t, w, h)
    R_(sl, l + 0.22, t + 0.255, 0.15, 0.15, fill=dot_color)
    T(sl, title, l + 0.50, t + 0.20, w - 0.70, 0.4, sz=title_sz, bold=True, color=TXT, anchor='t')
    PARA(sl, desc, l + 0.22, t + 0.72, w - 0.44, h - 0.9, sz=12, color=MUT, lh=1.22)


def stepcard(sl, l, t, w, h, num, color, title, desc, num_sz=40):
    card(sl, l, t, w, h)
    T(sl, str(num), l + 0.18, t + 0.08, 1.0, 0.62, sz=num_sz, bold=True, color=color, font=SERIF, anchor='t')
    T(sl, title, l + 0.20, t + 0.74, w - 0.40, 0.40, sz=13, bold=True, color=TXT, anchor='t', wrap=True)
    PARA(sl, desc, l + 0.20, t + 1.16, w - 0.40, h - 1.28, sz=11, color=MUT, lh=1.2)


def pbox(sl, l, t, w, h, color, title, caption=None, title_sz=14, fill=WH):
    card(sl, l, t, w, h, border=color, fill=fill, sh=False, lw=1.3)
    ty = t + (h * 0.28 if caption else h / 2 - 0.16)
    T(sl, title, l, ty, w, 0.36, sz=title_sz, bold=True, color=color, align=PP_ALIGN.CENTER, anchor='ctr', wrap=False)
    if caption:
        T(sl, caption, l + 0.05, t + h * 0.6, w - 0.1, 0.30, sz=9.5, color=MUT,
          align=PP_ALIGN.CENTER, anchor='ctr', wrap=False)


def pipeline(sl, l, t, box_w, box_h, gap, items, line_color=SUBT, title_sz=14):
    """items: lista de (titulo, cor, legenda-opcional)"""
    x = l
    for i, item in enumerate(items):
        title, color = item[0], item[1]
        cap = item[2] if len(item) > 2 else None
        pbox(sl, x, t, box_w, box_h, color, title, cap, title_sz=title_sz)
        if i < len(items) - 1:
            x1 = x + box_w + 0.04
            x2 = x + box_w + gap - 0.04
            LINE(sl, x1, t + box_h / 2, x2, t + box_h / 2, color=line_color, w=1.6)
        x += box_w + gap
    return x - gap


def make_table(sl, l, t, w, rows, col_fracs, header_h=0.46, row_h=0.46,
               sz=12, header_sz=12, mono_cols=(), bold_first_col=True,
               highlight_row=None, highlight_fill=C("FBEFE9")):
    nrows = len(rows); ncols = len(rows[0])
    total_h = header_h + (nrows - 1) * row_h
    gframe = sl.shapes.add_table(nrows, ncols, Inches(l), Inches(t), Inches(w), Inches(total_h))
    tbl = gframe.table
    tbl.first_row = False
    tbl.horz_banding = False
    for c in range(ncols):
        tbl.columns[c].width = Inches(w * col_fracs[c])
    tbl.rows[0].height = Inches(header_h)
    for r in range(1, nrows):
        tbl.rows[r].height = Inches(row_h)
    for r in range(nrows):
        for c in range(ncols):
            cell = tbl.cell(r, c)
            cell.margin_left = Inches(0.14); cell.margin_right = Inches(0.14)
            cell.margin_top = Inches(0.04); cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run(); run.text = rows[r][c]
            is_mono = (c in mono_cols) and r > 0
            run.font.name = MONO if is_mono else SANS
            run.font.size = Pt(header_sz if r == 0 else sz)
            run.font.bold = (r == 0) or (bold_first_col and c == 0 and r > 0)
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = TXT
                run.font.color.rgb = WH
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = highlight_fill if r == highlight_row else (ROWA if r % 2 else ROWB)
                run.font.color.rgb = TL if is_mono else (TXT if c == 0 else MUT)
    return tbl


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Capa
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bgfill(sl, CV)
sq_y = 1.50
for i in range(7):
    sq = R_(sl, 0.70 + i * 0.46, sq_y, 0.34, 0.34, fill=OR)
    sq.fill.fore_color.rgb = OR
    alpha_val = str(82000 - i * 11000)
    sf = sq.fill.fore_color._xFill.find(qn('a:srgbClr'))
    a_el = etree.SubElement(sf, qn('a:alpha')); a_el.set('val', alpha_val)

T(sl, "ESCOLA POLIT\u00c9CNICA \u00b7 PUCRS \u00b7 SISTEMAS DISTRIBU\u00cdDOS",
  1.80, 2.78, 10.5, 0.32, sz=12.5, bold=True, color=OR, font=SANS, spc=160, anchor='ctr')
T(sl, "RabbitMQ", 1.76, 3.10, 10.5, 1.10, sz=62, bold=True, color=WH, font=SERIF, anchor='t')
T(sl, "Mensageria com AMQP: arquitetura, garantias de entrega e ordena\u00e7\u00e3o",
  1.80, 4.28, 10.0, 0.50, sz=18, italic=True, color=CVSUB, font=SANS, anchor='t')
divider_h(sl, 1.80, 5.02, 6.80, color=CVDIV)
T(sl, "Arthur Real Sanchotene Ferreira \u00b7 Caroline da Rocha de Lima \u00b7 Giovani da Silva Cancherini",
  1.80, 5.14, 10.5, 0.34, sz=13, bold=True, color=WH, anchor='t')
T(sl, "Leonardo Vargas Schilling \u00b7 Osmar Sadi Nether Filho",
  1.80, 5.46, 10.5, 0.34, sz=13, color=CVMUT, anchor='t')
T(sl, "Prof. Fernando Dotti \u00b7 Porto Alegre, junho de 2026",
  1.80, 5.86, 10.5, 0.34, sz=12, italic=True, color=CVMUT, anchor='t')


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Agenda
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
head(sl, None, "Agenda", "Sete quest\u00f5es guiam nossa apresenta\u00e7\u00e3o, da motiva\u00e7\u00e3o at\u00e9 a demonstra\u00e7\u00e3o pr\u00e1tica.")
items = [
    "O que \u00e9 o RabbitMQ e por que existe",
    "Arquitetura e Funcionalidades",
    "Tipos de Falhas Suportadas",
    "Confiabilidade: modelos formais",
    "Ordena\u00e7\u00e3o de Mensagens",
    "Implementa\u00e7\u00e3o Interna",
    "Demonstra\u00e7\u00e3o Pr\u00e1tica: TicketLab",
]
row_h, gap = 0.555, 0.10
y = CONTENT_Y
for i, it in enumerate(items):
    R_(sl, L, y, 0.56, row_h, fill=OR)
    T(sl, "%02d" % (i + 1), L, y, 0.56, row_h, sz=17, bold=True, color=WH,
      font=SERIF, align=PP_ALIGN.CENTER, anchor='ctr')
    card(sl, L + 0.68, y, CW - 0.68, row_h)
    T(sl, it, L + 0.95, y, CW - 0.68 - 0.50, row_h, sz=14.5, bold=True, color=TXT, anchor='ctr')
    y += row_h + gap
pagenum(sl, 2)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — O Problema
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
head(sl, "CONTEXTO", "O Problema", "HTTP s\u00edncrono: uma cadeia que quebra pela parte mais fraca.")
T(sl, "Imagine o Lollapalooza abrindo vendas \u00e0s 10h: 50.000 pessoas clicando ao mesmo tempo.",
  L, CONTENT_Y, CW, 0.32, sz=12.5, italic=True, color=SUBT, anchor='t')

code_l, code_w = L, 6.55
code(sl, code_l, CONTENT_Y + 0.42, code_w, 1.55,
     "Usu\u00e1rio -> API -> Pagamento (aguarda 3s)\n"
     "              -> Estoque (aguarda 1s)\n"
     "              -> Email (aguarda 2s)\n"
     "         <- responde apos 6 segundos")
PARA(sl, [
    "Se Pagamento demorar 30s: timeout. Usu\u00e1rio v\u00ea erro 500.",
    "Se Estoque cair: toda a compra falha.",
    "2 usu\u00e1rios simult\u00e2neos: overbooking poss\u00edvel.",
], code_l, CONTENT_Y + 2.10, code_w, 0.95, sz=11.5, color=SUBT, lh=1.35)

right_l, right_w = 7.45, RM - 7.45
card(sl, right_l, CONTENT_Y + 0.42, right_w, 2.55)
T(sl, "O que falha", right_l + 0.24, CONTENT_Y + 0.62, right_w - 0.48, 0.32, sz=14, bold=True, color=OR, anchor='t')
bullets(sl, [
    "Timeout em cascata",
    "Falha se qualquer servi\u00e7o cair",
    "Overbooking",
    "Thread presa at\u00e9 o fim",
    "Escalar = problema",
], right_l + 0.24, CONTENT_Y + 1.05, right_w - 0.48, sz=12.5, dot_color=RD, text_color=TXT, dy=0.40)

ph_t = CONTENT_Y + 3.18
placeholder(sl, L, ph_t, CW, BOTTOM - ph_t, "v\u00eddeo opcional",
            "\"RabbitMQ in 100 Seconds\": Fireship (YouTube)",
            "Abrir no browser antes da apresenta\u00e7\u00e3o, reproduzir os primeiros 45s. Pausar antes do pr\u00f3ximo slide.")
pagenum(sl, 3)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — A Solu\u00e7\u00e3o
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
head(sl, "CONTEXTO", "A Solu\u00e7\u00e3o", "Desacoplar produ\u00e7\u00e3o de consumo via broker de mensagens.")

pipe_items = [
    ("Usu\u00e1rio", SUBT), ("API", TL), ("Fila", OR),
    ("worker 1", GN), ("worker 2", GN), ("worker 3", GN),
]
box_w, gap_w = 1.73, 0.22
pipe_t = CONTENT_Y + 0.10
pipeline(sl, L, pipe_t, box_w, 0.85, gap_w, pipe_items, title_sz=13.5)

code(sl, L, pipe_t + 1.05, 5.0, 0.42,
     '<- API responde imediato: { "status": "pending" }', sz=10.5, color=GN)

cards_t = pipe_t + 1.70
card_w = (CW - 2 * 0.24) / 3
dotcard(sl, L, cards_t, card_w, 1.55, GN, "Resposta imediata",
        ["API retorna na hora. Processamento ocorre em background."])
dotcard(sl, L + card_w + 0.24, cards_t, card_w, 1.55, OR, "Workers independentes",
        ["Cada etapa \u00e9 um servi\u00e7o separado. Falha num n\u00e3o para os outros."])
dotcard(sl, L + 2 * (card_w + 0.24), cards_t, card_w, 1.55, TL, "Escalabilidade natural",
        ["Adicionar workers = aumentar throughput sem mudar a API."])
pagenum(sl, 4)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — O que \u00e9 o RabbitMQ?
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
head(sl, "VIS\u00c3O GERAL", "O que \u00e9 o RabbitMQ?", "Message broker open-source: o intermedi\u00e1rio confi\u00e1vel.")
y = CONTENT_Y
sections = [
    ("Message Broker", [
        "Intermedi\u00e1rio entre quem envia e quem recebe mensagens.",
        "Desacopla produtor de consumidor no tempo e no espa\u00e7o. N\u00e3o \u00e9 um banco de dados: \u00e9 um canal confi\u00e1vel.",
    ]),
    ("Protocolo AMQP", [
        "Implementa AMQP 0-9-1 (Advanced Message Queuing Protocol), protocolo aberto e bin\u00e1rio com sem\u00e2ntica bem definida de ACK.",
        "Tamb\u00e9m suporta STOMP, MQTT e AMQP 1.0 via plugins.",
    ]),
    ("Erlang / OTP", [
        "Escrito em Erlang, linguagem criada para sistemas de telecom com 99,9999999% de disponibilidade (nine nines).",
        "Isolamento de processos e recupera\u00e7\u00e3o de falhas nativos.",
    ]),
]
for i, (title, desc) in enumerate(sections):
    text_block(sl, title, desc, L, y, CW, title_color=OR, title_sz=15, body_sz=12.5)
    y += 1.30
    if i < len(sections) - 1:
        divider_h(sl, L, y - 0.14, CW)
pagenum(sl, 5)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Onde \u00e9 utilizado?
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
head(sl, "VIS\u00c3O GERAL", "Onde \u00e9 utilizado?", "Qualquer sistema que precisa desacoplar e escalar.")
cases = [
    ("E-commerce", ["Checkout, pagamento, estoque e notifica\u00e7\u00e3o.", "Exatamente o nosso projeto TicketLab."], OR),
    ("Sistemas banc\u00e1rios", ["Processamento ass\u00edncrono de transa\u00e7\u00f5es.", "Auditoria e filas de concilia\u00e7\u00e3o."], TL),
    ("IoT e sensores", ["Ingest\u00e3o de eventos de milhares de dispositivos.", "Fanout para m\u00faltiplos consumidores anal\u00edticos."], AM),
    ("Microservi\u00e7os", ["Comunica\u00e7\u00e3o entre bounded contexts sem acoplamento temporal.", "Event-driven architecture."], GN),
    ("Reservas / Ingressos", ["Controle de concorr\u00eancia em estoque limitado.", "Fila garante um worker por vez no estoque."], OR),
]
n = len(cases)
cw = (CW - (n - 1) * 0.20) / n
ch = 3.85
for i, (title, desc, accent) in enumerate(cases):
    cx = L + i * (cw + 0.20)
    card(sl, cx, CONTENT_Y, cw, ch)
    R_(sl, cx, CONTENT_Y, cw, 0.07, fill=accent)
    T(sl, title, cx + 0.18, CONTENT_Y + 0.26, cw - 0.36, 0.65, sz=13.5, bold=True, color=TXT, wrap=True, anchor='t')
    PARA(sl, desc, cx + 0.18, CONTENT_Y + 1.00, cw - 0.36, ch - 1.15, sz=11, color=MUT, lh=1.28)
pagenum(sl, 6)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Arquitetura Interna
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
head(sl, "ARQUITETURA", "Arquitetura Interna", "Producer \u2192 Exchange \u2192 Queue \u2192 Consumer.", badge_color=TL)

img_t = CONTENT_Y + 0.15
card(sl, L, img_t, CW, 2.15)
picture_fit(sl, os.path.join(ASSETS, "architecture-flow.png"), L + 0.30, img_t + 0.20, CW - 0.60, 1.75)

comp_t = img_t + 2.15 + 0.28
comps = [
    ("Producer", OR, "Envia mensagens ao Exchange. N\u00e3o conhece as filas."),
    ("Exchange", TL, "Recebe e roteia via binding + routing key. N\u00e3o armazena."),
    ("Queue", AM, "Armazena mensagens at\u00e9 o consumer estar pronto."),
    ("Consumer", GN, "Processa e envia ACK. S\u00f3 ent\u00e3o o broker remove da fila."),
]
cw4 = (CW - 3 * 0.20) / 4
for i, (title, color, desc) in enumerate(comps):
    cx = L + i * (cw4 + 0.20)
    dotcard(sl, cx, comp_t, cw4, BOTTOM - comp_t, color, title, [desc], title_sz=13.5)
pagenum(sl, 7)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Exchange Direct
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
head(sl, "ARQUITETURA", "Exchange Direct", "Roteamento pela chave exata: cada mensagem vai para uma \u00fanica fila.", badge_color=TL)

left_w = 5.95
diag_t = CONTENT_Y + 0.30
card(sl, L, diag_t, left_w, 3.35)
pbox(sl, L + 0.28, diag_t + 0.30, 1.55, 0.62, SUBT, "Producer", title_sz=12)
T(sl, "key=\"payment\"", L + 0.10, diag_t + 1.00, 1.95, 0.26, sz=9.5, italic=True, color=MUT, align=PP_ALIGN.CENTER)
ex_cx, ex_cy, ex_d = L + 2.62, diag_t + 0.61, 0.85
OVAL(sl, ex_cx, ex_cy - ex_d / 2, ex_d, ex_d, fill=TL)
T(sl, "Direct\nExchange", ex_cx, ex_cy - 0.26, ex_d, 0.52, sz=10, bold=True, color=WH,
  align=PP_ALIGN.CENTER, anchor='ctr', wrap=True)
LINE(sl, L + 0.28 + 1.55 + 0.04, diag_t + 0.61, ex_cx - 0.04, ex_cy, color=SUBT, w=1.6)

q_l = ex_cx + ex_d + 0.55
q_w, q_h, q_gap = 2.30, 0.56, 0.20
queues = [("payment-queue", True), ("stock-queue", False), ("notification-queue", False)]
qy = diag_t + 0.15
for name, match in queues:
    col = GN if match else BD
    txtcol = TXT if match else MUT
    pbox(sl, q_l, qy, q_w, q_h, col if match else MUT, name, title_sz=10.5)
    LINE(sl, ex_cx + ex_d / 2 + 0.02, ex_cy, q_l - 0.04, qy + q_h / 2,
         color=(GN if match else BD), w=(2.0 if match else 1.0))
    if match:
        T(sl, "match!", q_l + q_w + 0.10, qy + q_h / 2 - 0.13, 0.9, 0.26, sz=10, bold=True, color=GN)
    qy += q_h + q_gap
PARA(sl, ["Apenas a routing key id\u00eantica recebe a mensagem; as demais filas n\u00e3o s\u00e3o tocadas."],
     L + 0.28, diag_t + 2.65, left_w - 0.56, 0.55, sz=10.5, italic=True, color=MUT, lh=1.25)

right_l = L + left_w + 0.30
right_w = RM - right_l
T(sl, "No TicketLab", right_l, CONTENT_Y, right_w, 0.32, sz=14, bold=True, color=OR, anchor='t')
code(sl, right_l, CONTENT_Y + 0.40, right_w, 1.35,
     "ticket-sales -> payment\n  -> payment-queue\n"
     "payment-approved -> stock\n  -> stock-queue\n"
     "stock-confirmed -> notification\n  -> notification-queue", sz=10)
T(sl, "Por que Direct?", right_l, CONTENT_Y + 1.95, right_w, 0.32, sz=14, bold=True, color=OR, anchor='t')
bullets(sl, [
    "Pipeline linear: cada etapa tem uma fila dedicada",
    "Routing simples e previs\u00edvel",
    "F\u00e1cil de monitorar no Management UI",
], right_l, CONTENT_Y + 2.40, right_w, sz=12, dot_color=TL, text_color=TXT, dy=0.46)
pagenum(sl, 8)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Outros Tipos de Exchange
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
head(sl, "ARQUITETURA", "Outros Tipos de Exchange", "Fanout \u00b7 Topic \u00b7 Headers.", badge_color=TL)
others = [
    ("Fanout", ["Copia para TODAS as filas vinculadas. Ignora routing key.",
                "Uso: broadcast de eventos (logs, auditoria, notifica\u00e7\u00f5es globais)."],
     "fanout-exchange.png"),
    ("Topic", ["Padr\u00e3o com wildcards: * (uma palavra) e # (zero ou mais).",
               "\"order.#\" captura \"order.created\", \"order.paid\"... Uso: sistemas de eventos flex\u00edveis."],
     "topic-exchange.png"),
    ("Headers", ["Roteia por atributos do cabe\u00e7alho AMQP.",
                 "Mais expressivo que Topic, por\u00e9m mais raro na pr\u00e1tica. Uso: roteamento por metadados arbitr\u00e1rios."],
     "rabbitmq-headers-exchange.png"),
]
n = len(others); cw = (CW - (n - 1) * 0.24) / n
img_h = 2.55
for i, (title, desc, img) in enumerate(others):
    cx = L + i * (cw + 0.24)
    card(sl, cx, CONTENT_Y, cw, 1.62)
    T(sl, title, cx + 0.22, CONTENT_Y + 0.16, cw - 0.44, 0.38, sz=18, bold=True, color=OR, font=SERIF, anchor='t')
    PARA(sl, desc, cx + 0.22, CONTENT_Y + 0.58, cw - 0.44, 1.0, sz=10.5, color=TXT, lh=1.2)
    img_t = CONTENT_Y + 1.62 + 0.16
    card(sl, cx, img_t, cw, img_h)
    picture_fit(sl, os.path.join(ASSETS, img), cx + 0.10, img_t + 0.10, cw - 0.20, img_h - 0.20)
pagenum(sl, 9)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Funcionalidades
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
head(sl, "ARQUITETURA", "Funcionalidades", "O que o RabbitMQ oferece al\u00e9m do roteamento b\u00e1sico.", badge_color=TL)
feats = [
    ("Mensagens Persistentes", "delivery_mode=PERSISTENT grava em disco antes de confirmar. Sem isso, reinicio apaga tudo.", OR),
    ("Filas Durable", "Metadados da fila ficam no disco. Broker reinicia: fila existe e mensagens voltam.", TL),
    ("ACK Manual", "Consumer confirma s\u00f3 ap\u00f3s processar e salvar. Broker s\u00f3 remove ao receber ACK.", AM),
    ("Dead Letter Exchange", "NACK sem requeue envia para DLX: dead letter queue para an\u00e1lise e reprocessamento.", GN),
    ("QoS / Prefetch", "prefetch_count limita mensagens em voo por consumer. Backpressure autom\u00e1tico.", OR),
    ("Management UI / API", "Interface web + HTTP API para monitorar filas, consumers e publicar mensagens.", TL),
]
cols, rows = 3, 2
cw3 = (CW - (cols - 1) * 0.22) / cols
ch2 = 2.05
for i, (title, desc, color) in enumerate(feats):
    r, c = divmod(i, cols)
    cx = L + c * (cw3 + 0.22)
    cy = CONTENT_Y + r * (ch2 + 0.18)
    dotcard(sl, cx, cy, cw3, ch2, color, title, [desc], title_sz=13)
pagenum(sl, 10)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Falha 1: Crash-stop
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
head(sl, "FALHAS", "Falha 1: Crash-stop do Worker", "O container morre; nenhuma mensagem se perde.", badge_color=RD)

dl, dw = L, 7.15
row1_t = CONTENT_Y
bw, bh = 2.05, 0.78
pbox(sl, dl, row1_t, bw, bh, SUBT, "Fila", "payment-queue", title_sz=13)
pbox(sl, dl + bw + 1.00, row1_t, bw, bh, GN, "Worker A", "processando", title_sz=13)
pbox(sl, dl + (bw + 1.00) * 2, row1_t, bw, bh, GN, "Worker B", "aguardando", title_sz=13)
LINE(sl, dl + bw + 0.04, row1_t + bh / 2, dl + bw + 0.96, row1_t + bh / 2, color=SUBT, w=1.6)
LINE(sl, dl + (bw + 1.00) + bw + 0.04, row1_t + bh / 2, dl + (bw + 1.00) * 2 + 0.04, row1_t + bh * 0.5, color=SUBT, w=1.6)

row2_t = row1_t + bh + 0.55
ax = dl + bw + 1.00 + bw / 2
T(sl, "\u2717 Worker A cai (SIGKILL)", dl + bw + 1.00, row2_t - 0.40, bw, 0.30, sz=11.5, bold=True, color=RD,
  align=PP_ALIGN.CENTER, anchor='ctr')
LINE(sl, ax, row2_t - 0.06, ax, row2_t + 0.46, color=RD, w=1.8)
box2_w = 2.55
pbox(sl, dl, row2_t + 0.55, box2_w, bh, AM, "msgs unacked", "voltam para ready", title_sz=12.5)
pbox(sl, dl + box2_w + 1.00, row2_t + 0.55, box2_w, bh, GN, "Worker B assume", None, title_sz=12.5)
LINE(sl, dl + box2_w + 0.04, row2_t + 0.55 + bh / 2, dl + box2_w + 0.96, row2_t + 0.55 + bh / 2, color=SUBT, w=1.6)

rl = dl + dw + 0.30
rw = RM - rl
card(sl, rl, CONTENT_Y, rw, 2.85)
T(sl, "Por que n\u00e3o perde mensagens?", rl + 0.22, CONTENT_Y + 0.18, rw - 0.44, 0.55, sz=13.5, bold=True, color=TXT, wrap=True, anchor='t')
bullets(sl, [
    "TCP disconnect detectado: RabbitMQ libera msgs unacked",
    "Outras msgs na fila ficam em estado ready, esperando",
    "Worker B (ou C) pega as mensagens liberadas",
    "Sem interven\u00e7\u00e3o manual: \u00e9 autom\u00e1tico",
], rl + 0.22, CONTENT_Y + 0.85, rw - 0.44, sz=11.5, dot_color=GN, text_color=TXT, dy=0.46)

bot_t = row2_t + 0.55 + bh + 0.22
code(sl, dl, bot_t, dw, BOTTOM - bot_t,
     "make send-1000        # enche a fila\n"
     "make kill-payment      # mata um worker abruptamente\n"
     "make logs-payment      # outro worker assume as msgs", sz=10)
placeholder(sl, rl, bot_t, rw, BOTTOM - bot_t, "gif", "Terminal: kill-payment + logs-payment")
pagenum(sl, 11)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Falha 2: Crash-recovery + Retry + DLQ
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
head(sl, "FALHAS", "Falha 2: Crash-recovery + Retry + DLQ",
     "O broker reinicia \u00b7 a l\u00f3gica falha \u00b7 o destino final.", badge_color=RD, title_sz=27)

half_w = (CW - 0.30) / 2
card(sl, L, CONTENT_Y, half_w, BOTTOM - CONTENT_Y)
T(sl, "Crash-recovery do Broker", L + 0.22, CONTENT_Y + 0.18, half_w - 0.44, 0.34, sz=14, bold=True, color=OR, anchor='t')
bullets(sl, [
    "RabbitMQ reinicia (docker restart)",
    "Mensagens PERSISTENT + filas durable ficam no disco",
    "Workers reconectam automaticamente (connect_robust)",
    "Sistema retoma de onde parou",
], L + 0.22, CONTENT_Y + 0.62, half_w - 0.44, sz=11.5, dot_color=OR, text_color=TXT, dy=0.46)
code(sl, L + 0.22, CONTENT_Y + 2.55, half_w - 0.44, 0.40, "make restart-rabbit", sz=10.5)
placeholder(sl, L + 0.22, CONTENT_Y + 3.10, half_w - 0.44, BOTTOM - (CONTENT_Y + 3.10) - 0.18, "screenshot",
            "RabbitMQ UI ap\u00f3s restart: filas com msgs intactas")

rl = L + half_w + 0.30
card(sl, rl, CONTENT_Y, half_w, BOTTOM - CONTENT_Y)
T(sl, "Retry com Dead Letter Queue", rl + 0.22, CONTENT_Y + 0.18, half_w - 0.44, 0.34, sz=14, bold=True, color=TL, anchor='t')
code(sl, rl + 0.22, CONTENT_Y + 0.60, half_w - 0.44, 1.95,
     "msg -> worker -> falha (RuntimeError)\n"
     "  -> republica x-retry-count: 1\n"
     "  -> falha de novo: x-retry-count: 2\n"
     "  -> falha de novo: x-retry-count: 3\n"
     "  -> max_retries atingido!\n"
     "  -> NACK(requeue=False) -> DLX -> DLQ", sz=10.5)
placeholder(sl, rl + 0.22, CONTENT_Y + 2.75, half_w - 0.44, BOTTOM - (CONTENT_Y + 2.75) - 0.18, "screenshot",
            "RabbitMQ UI: dead-letter-queue com N mensagens")
pagenum(sl, 12)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Modelos de Falha: Resumo
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
head(sl, "FALHAS", "Modelos de Falha: Resumo", "O que suportamos vs. o que n\u00e3o suportamos.", badge_color=RD)
T(sl, "\u2713 SUPORTADO", L, CONTENT_Y, CW, 0.28, sz=11.5, bold=True, color=GN, spc=120)
sup = [
    ("Crash-stop", "Worker morre: mensagens voltam para a fila automaticamente."),
    ("Crash-recovery", "Broker reinicia: mensagens PERSISTENT sobrevivem no disco."),
    ("Omiss\u00e3o com retry", "Falhas de l\u00f3gica geram DLQ ap\u00f3s N tentativas."),
    ("Parti\u00e7\u00e3o de rede", "connect_robust reconecta com backoff exponencial."),
]
cw2 = (CW - 0.22) / 2
ch1 = 1.18
y0 = CONTENT_Y + 0.36
for i, (title, desc) in enumerate(sup):
    r, c = divmod(i, 2)
    cx = L + c * (cw2 + 0.22)
    cy = y0 + r * (ch1 + 0.14)
    dotcard(sl, cx, cy, cw2, ch1, GN, title, [desc], title_sz=13)

y1 = y0 + 2 * (ch1 + 0.14) + 0.20
T(sl, "\u2717 N\u00c3O SUPORTADO", L, y1, CW, 0.28, sz=11.5, bold=True, color=RD, spc=120)
nsup = [
    ("Falhas Byzantinas", "N\u00f3 malicioso enviando mensagens forjadas ou corrompidas."),
    ("Exatamente-uma-vez nativo", "At-least-once: duplicatas poss\u00edveis (tratamos com idempot\u00eancia)."),
]
y2 = y1 + 0.36
for i, (title, desc) in enumerate(nsup):
    cx = L + i * (cw2 + 0.22)
    dotcard(sl, cx, y2, cw2, ch1, RD, title, [desc], title_sz=13)
T(sl, "Falhas Byzantinas requerem protocolo BFT dedicado (PBFT, Tendermint): fora do escopo do RabbitMQ.",
  L, y2 + ch1 + 0.12, CW, 0.30, sz=10.5, italic=True, color=MUT)
pagenum(sl, 13)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Confiabilidade: Fair-Loss vs Perfect Link
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
head(sl, "CONFIABILIDADE", "Confiabilidade: Fair-Loss vs Perfect Link",
     "O que acontece com as mensagens quando o broker cai?", badge_color=TL, title_sz=26)
half_w = (CW - 0.30) / 2
card(sl, L, CONTENT_Y, half_w, BOTTOM - CONTENT_Y)
T(sl, "Fair-Loss Link", L + 0.22, CONTENT_Y + 0.18, half_w - 0.44, 0.34, sz=14, bold=True, color=AM, anchor='t')
PARA(sl, ["Mensagens podem se perder, mas n\u00e3o s\u00e3o inventadas."], L + 0.22, CONTENT_Y + 0.58, half_w - 0.44, 0.4, sz=12, color=TXT)
code(sl, L + 0.22, CONTENT_Y + 1.00, half_w - 0.44, 0.85,
     "exchange -> non-durable\nfila -> non-durable\ndelivery -> NON_PERSISTENT (RAM)", sz=10.5)
T(sl, "Resultado:", L + 0.22, CONTENT_Y + 1.98, half_w - 0.44, 0.3, sz=12, bold=True, color=TXT)
bullets(sl, [
    "Broker reinicia: mensagens na RAM somem",
    "Nenhum aviso ao producer",
    "Garantia? Nenhuma.",
], L + 0.22, CONTENT_Y + 2.32, half_w - 0.44, sz=11.5, dot_color=AM, text_color=TXT, dy=0.40)

rl = L + half_w + 0.30
card(sl, rl, CONTENT_Y, half_w, BOTTOM - CONTENT_Y)
T(sl, "Perfect Link (at-least-once)", rl + 0.22, CONTENT_Y + 0.18, half_w - 0.44, 0.34, sz=14, bold=True, color=GN, anchor='t')
PARA(sl, ["Toda mensagem enviada \u00e9 eventualmente entregue."], rl + 0.22, CONTENT_Y + 0.58, half_w - 0.44, 0.4, sz=12, color=TXT)
code(sl, rl + 0.22, CONTENT_Y + 1.00, half_w - 0.44, 1.10,
     "exchange -> durable\nfila -> durable\ndelivery -> PERSISTENT (disco)\nACK -> manual\nconexao -> connect_robust", sz=10.5)
T(sl, "Ressalva: at-least-once", rl + 0.22, CONTENT_Y + 2.22, half_w - 0.44, 0.3, sz=12, bold=True, color=TXT)
PARA(sl, ["Worker processa e falha antes do ACK: recebe a mensagem de novo. Idempot\u00eancia (pr\u00f3ximos slides) trata isso."],
     rl + 0.22, CONTENT_Y + 2.56, half_w - 0.44, 0.85, sz=11.5, color=MUT, lh=1.3)
pagenum(sl, 14)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — BEB vs Reliable Broadcast
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
head(sl, "CONFIABILIDADE", "BEB vs Reliable Broadcast", "O producer falha no meio: o que acontece?", badge_color=TL)
half_w = (CW - 0.30) / 2
card(sl, L, CONTENT_Y, half_w, BOTTOM - CONTENT_Y)
T(sl, "Best Effort Broadcast", L + 0.22, CONTENT_Y + 0.18, half_w - 0.44, 0.34, sz=14, bold=True, color=AM, anchor='t')
PARA(sl, ["O emissor dispara e n\u00e3o garante se ele mesmo falhar."], L + 0.22, CONTENT_Y + 0.58, half_w - 0.44, 0.35, sz=12, color=TXT)
code(sl, L + 0.22, CONTENT_Y + 1.00, half_w - 0.44, 0.95,
     "# publisher.py - sem publisher confirms\n"
     "await exchange.publish(msg, key='payment')\n"
     "# broker cai aqui -> msg perdida", sz=10)
bullets(sl, [
    "Enviamos 1000 mensagens",
    "Broker cai na mensagem 500",
    "Mensagens 501 a 1000 perdidas",
    "Producer n\u00e3o tem como saber",
], L + 0.22, CONTENT_Y + 2.10, half_w - 0.44, sz=11.5, dot_color=AM, text_color=TXT, dy=0.38)

rl = L + half_w + 0.30
card(sl, rl, CONTENT_Y, half_w, BOTTOM - CONTENT_Y)
T(sl, "Reliable Broadcast (aprox.)", rl + 0.22, CONTENT_Y + 0.18, half_w - 0.44, 0.34, sz=14, bold=True, color=GN, anchor='t')
PARA(sl, ["Se um processo correto entrega m, todos entregam m."], rl + 0.22, CONTENT_Y + 0.58, half_w - 0.44, 0.35, sz=12, color=TXT)
bullets(sl, [
    "Mensagem com PERSISTENT nunca some",
    "ACK manual: broker guarda at\u00e9 confirmar",
    "Worker ativo processa eventualmente",
], rl + 0.22, CONTENT_Y + 1.00, half_w - 0.44, sz=11.5, dot_color=GN, text_color=TXT, dy=0.40)
T(sl, "O que ainda falta para RB formal:", rl + 0.22, CONTENT_Y + 2.45, half_w - 0.44, 0.3, sz=12, bold=True, color=TXT)
bullets(sl, [
    "Publisher confirms (sem eles = BEB)",
    "Se producer cai antes de publicar, n\u00e3o h\u00e1 recupera\u00e7\u00e3o",
], rl + 0.22, CONTENT_Y + 2.78, half_w - 0.44, sz=11, dot_color=GN, text_color=MUT, dy=0.36)
pagenum(sl, 15)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Uniform Reliable Broadcast
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
head(sl, "CONFIABILIDADE", "Uniform Reliable Broadcast", "O n\u00edvel mais alto: como chegamos perto.", badge_color=TL)
T(sl, "Defini\u00e7\u00e3o de URB: se qualquer processo (correto ou com falha) entrega m, TODOS entregam m. Requer consenso.",
  L, CONTENT_Y, CW, 0.45, sz=12.5, italic=True, color=SUBT, wrap=True)

half_w = (CW - 0.30) / 2
card(sl, L, CONTENT_Y + 0.55, half_w, 2.55, border=RD)
T(sl, "\u2717 RabbitMQ n\u00e3o garante URB nativo", L + 0.22, CONTENT_Y + 0.73, half_w - 0.44, 0.34, sz=13.5, bold=True, color=RD, anchor='t', wrap=True)
PARA(sl, [
    "Worker A: recebe a mensagem, processa parcialmente, salva metade no banco e falha antes do ACK.",
    "Worker B: recebe a mesma mensagem de novo. O estado parcial do Worker A j\u00e1 existe: sem consenso sobre quem entregou.",
], L + 0.22, CONTENT_Y + 1.18, half_w - 0.44, 1.85, sz=11.5, color=TXT, lh=1.3)

rl = L + half_w + 0.30
card(sl, rl, CONTENT_Y + 0.55, half_w, 2.55, border=GN)
T(sl, "\u2713 Nossa solu\u00e7\u00e3o: idempot\u00eancia", rl + 0.22, CONTENT_Y + 0.73, half_w - 0.44, 0.34, sz=13.5, bold=True, color=GN, anchor='t')
code(sl, rl + 0.22, CONTENT_Y + 1.15, half_w - 0.44, 1.20,
     "# notification-worker/worker.py\n"
     "if order.status == 'confirmed':\n"
     "    return  # ja processado, ignora", sz=10.5)
PARA(sl, ["Resultado: exactly-once sem\u00e2ntico para o usu\u00e1rio, mesmo com at-least-once na fila."],
     rl + 0.22, CONTENT_Y + 2.42, half_w - 0.44, 0.6, sz=11.5, color=TXT, lh=1.3)

T(sl, "Quorum Queues (RabbitMQ 3.8+ com Raft) aproximam o URB nativamente: pr\u00f3ximo passo real.",
  L, CONTENT_Y + 3.25, CW, 0.32, sz=11, italic=True, color=MUT)
pagenum(sl, 16)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Comparativo: os 5 Modelos
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
head(sl, "CONFIABILIDADE", "Comparativo: os 5 Modelos", "De Fair-Loss at\u00e9 URB: onde o RabbitMQ se encaixa.", badge_color=TL)
rows = [
    ["Modelo", "Perde msg?", "Duplicatas?", "Crash producer", "No projeto"],
    ["Fair-Loss Link", "Sim", "N\u00e3o", "n/a", "Sem persist\u00eancia"],
    ["Perfect Link", "N\u00e3o", "Poss\u00edvel", "Perde msgs", "PERSISTENT + ACK + durable"],
    ["BEB", "Poss\u00edvel", "N\u00e3o", "Perde msgs", "POST /batch (sem confirms)"],
    ["Reliable Broadcast", "N\u00e3o", "Poss\u00edvel", "Perde msgs", "Aprox. com persist\u00eancia"],
    ["Uniform RB", "N\u00e3o", "N\u00e3o", "Nenhuma", "N\u00e3o nativo (idempot\u00eancia)"],
]
make_table(sl, L, CONTENT_Y, CW, rows, col_fracs=[0.20, 0.15, 0.15, 0.18, 0.32],
           header_h=0.48, row_h=0.50, sz=12, header_sz=12.5, highlight_row=2)
T(sl, "Linha destacada = nossa implementa\u00e7\u00e3o principal.", L, CONTENT_Y + 0.48 + 5 * 0.50 + 0.18, CW, 0.3,
  sz=11, italic=True, color=MUT)
pagenum(sl, 17)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Ordenação: FIFO com 1 Consumer
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
head(sl, "ORDENA\u00c7\u00c3O", "Ordena\u00e7\u00e3o: FIFO com 1 Consumer",
     "A fila garante a ordem. Um consumidor a respeita.", badge_color=AM)

row_t = CONTENT_Y + 0.10
pipe_items = [("msg 1", SUBT), ("msg 2", SUBT), ("msg 3", SUBT), ("msg 4", SUBT),
              ("msg 5", SUBT), ("Worker", GN, "1 consumer")]
box_w = (CW - 5 * 0.16) / 6
pipeline(sl, L, row_t, box_w, 0.85, 0.16, pipe_items, title_sz=13)

out_t = row_t + 0.85 + 0.25
code(sl, L, out_t, CW, 0.42,
     "ACK por mensagem, sempre na mesma ordem:  msg1 -> ack -> msg2 -> ack -> msg3 -> ack -> msg4 -> ack -> msg5 -> ack",
     sz=10.5, color=GN, fill=C("EAF5EE"))

card_t = out_t + 0.42 + 0.25
card(sl, L, card_t, CW, 1.35, border=AM)
T(sl, "Propriedade FIFO garantida pelo RabbitMQ", L + 0.22, card_t + 0.16, CW - 0.44, 0.32, sz=13, bold=True, color=AM)
PARA(sl, ["A especifica\u00e7\u00e3o AMQP 0-9-1 define que mensagens publicadas com a mesma routing key e priority "
          "chegam \u00e0 fila na ordem em que foram publicadas, e s\u00e3o entregues nessa ordem a um \u00fanico consumer."],
     L + 0.22, card_t + 0.54, CW - 0.44, 0.75, sz=11.5, color=TXT, lh=1.3)

warn_t = card_t + 1.35 + 0.22
R_(sl, L, warn_t + 0.065, 0.09, 0.09, fill=RD)
T2(sl, [{'text': 'Mas: ', 'bold': True, 'color': RD, 'sz': 12.5},
        {'text': 'com N consumers em paralelo, a ordem global de processamento N\u00c3O \u00e9 garantida (pr\u00f3ximo slide).',
         'color': TXT, 'sz': 12.5}], L + 0.26, warn_t, CW - 0.26, 0.4)
pagenum(sl, 18)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Ordenação: N Consumers
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
head(sl, "ORDENA\u00c7\u00c3O", "Ordena\u00e7\u00e3o: N Consumers", "Paralelo quebra a ordem global. Como resolver.", badge_color=AM)

row_t = CONTENT_Y + 0.05
col_gap = 0.28
col_w = (CW - 2 * col_gap) / 3
col_h = 2.05
c1, c2, c3 = L, L + col_w + col_gap, L + 2 * (col_w + col_gap)
mid_y = row_t + col_h / 2

card(sl, c1, row_t, col_w, col_h)
T(sl, "Fila", c1 + 0.20, row_t + 0.16, col_w - 0.40, 0.30, sz=13, bold=True, color=TL)
msgs = ["msg 1", "msg 2", "msg 3", "msg 4"]
mb_w, mb_h, mb_gap = 0.74, 0.52, 0.08
mx0 = c1 + (col_w - (4 * mb_w + 3 * mb_gap)) / 2
for i, m in enumerate(msgs):
    pbox(sl, mx0 + i * (mb_w + mb_gap), row_t + 1.20, mb_w, mb_h, SUBT, m, title_sz=9.5)
T(sl, "ordem de chegada: 1, 2, 3, 4", c1 + 0.20, row_t + 1.78, col_w - 0.40, 0.26, sz=9.5, italic=True, color=MUT)

card(sl, c2, row_t, col_w, col_h)
T(sl, "Workers", c2 + 0.20, row_t + 0.16, col_w - 0.40, 0.30, sz=13, bold=True, color=TL)
pbox(sl, c2 + 0.20, row_t + 0.58, col_w - 0.40, 0.62, SUBT, "Worker A", "lento \u00b7 msg 1", title_sz=12)
pbox(sl, c2 + 0.20, row_t + 1.32, col_w - 0.40, 0.62, GN, "Worker B", "r\u00e1pido \u00b7 msg 2", title_sz=12)

card(sl, c3, row_t, col_w, col_h)
T(sl, "Resultado", c3 + 0.20, row_t + 0.16, col_w - 0.40, 0.30, sz=13, bold=True, color=RD)
T2(sl, [{'text': 'Chegada: ', 'bold': True, 'sz': 12, 'color': TXT}, {'text': '1, 2, 3, 4', 'sz': 12, 'color': MUT, 'font': MONO}],
   c3 + 0.20, row_t + 0.62, col_w - 0.40, 0.30)
T2(sl, [{'text': 'Processado: ', 'bold': True, 'sz': 12, 'color': TXT}, {'text': '2, 4, 1, 3', 'sz': 12, 'color': RD, 'font': MONO}],
   c3 + 0.20, row_t + 0.96, col_w - 0.40, 0.30)
R_(sl, c3 + 0.20, row_t + 1.42, 0.09, 0.09, fill=RD)
T(sl, "Ordem global quebrada", c3 + 0.20 + 0.22, row_t + 1.355, col_w - 0.62, 0.30, sz=11.5, bold=True, color=RD)

LINE(sl, c1 + col_w + 0.04, mid_y, c2 - 0.04, mid_y, color=SUBT, w=1.6)
LINE(sl, c2 + col_w + 0.04, mid_y, c3 - 0.04, mid_y, color=SUBT, w=1.6)

lower_t = row_t + col_h + 0.30
half_w = (CW - 0.30) / 2
code(sl, L, lower_t, half_w, BOTTOM - lower_t,
     "make scale-workers N=5 && make send-100\n"
     "make db-query\n"
     "# processed_at N\u00c3O respeita created_at", sz=10.5)

rl = L + half_w + 0.30
card(sl, rl, lower_t, half_w, BOTTOM - lower_t)
T(sl, "Como garantir ordem total, se necess\u00e1rio", rl + 0.22, lower_t + 0.16, half_w - 0.44, 0.55, sz=13, bold=True, color=TXT, wrap=True)
text_block(sl, "Particionamento", ["1 worker por event_id, como Kafka faz com partition keys. Simples e eficaz."],
           rl + 0.22, lower_t + 0.62, half_w - 0.44, title_color=OR, title_sz=12, body_sz=10.5)
text_block(sl, "Sequencer", ["N\u00f3 central atribui n\u00famero de sequ\u00eancia antes de publicar. Gargalo, mas garante ordem."],
           rl + 0.22, lower_t + 1.30, half_w - 0.44, title_color=TL, title_sz=12, body_sz=10.5)
text_block(sl, "Kafka", ["Nativo: cada parti\u00e7\u00e3o \u00e9 uma fila FIFO. Workers de uma parti\u00e7\u00e3o respeitam a ordem."],
           rl + 0.22, lower_t + 1.98, half_w - 0.44, title_color=GN, title_sz=12, body_sz=10.5)
pagenum(sl, 19)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Implementação: Persistência + ACK
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
head(sl, "IMPLEMENTA\u00c7\u00c3O", "Implementa\u00e7\u00e3o: Persist\u00eancia + ACK",
     "Como o RabbitMQ garante que mensagens n\u00e3o somem.", badge_color=GN)

steps = [
    (1, OR, "Publisher", ["publish(PERSISTENT)", "Grava em disco antes de retornar."]),
    (2, TL, "Disco", ["mnesia / msg store", "Arquivo em /var/lib/rabbitmq."]),
    (3, AM, "Consumer", ["processa()", "L\u00f3gica de neg\u00f3cio + commit no banco."]),
    (4, GN, "ACK", ["channel.ack()", "Broker remove da fila (disco)."]),
]
cw4 = (CW - 3 * 0.20) / 4
step_h = 2.15
step_t = CONTENT_Y
for i, (num, color, title, desc) in enumerate(steps):
    cx = L + i * (cw4 + 0.20)
    stepcard(sl, cx, step_t, cw4, step_h, num, color, title, desc, num_sz=36)
    if i < len(steps) - 1:
        ay = step_t + 0.42
        LINE(sl, cx + cw4 + 0.04, ay, cx + cw4 + 0.20 - 0.04, ay, color=SUBT, w=1.6)

warn_t = step_t + step_h + 0.28
card(sl, L, warn_t, CW, 1.05, border=RD)
T(sl, "\u26a0 Sem ACK", L + 0.24, warn_t + 0.16, 2.0, 0.32, sz=13, bold=True, color=RD)
PARA(sl, ["Crash antes do ACK: a mensagem volta para o estado ready e \u00e9 entregue de novo. "
          "Isso \u00e9 o pre\u00e7o do at-least-once: necess\u00e1rio idempot\u00eancia no processamento."],
     L + 0.24, warn_t + 0.52, CW - 0.48, 0.45, sz=11.5, color=TXT, lh=1.25)
pagenum(sl, 20)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Implementação: Anti-Overbooking
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
head(sl, "IMPLEMENTA\u00c7\u00c3O", "Implementa\u00e7\u00e3o: Anti-Overbooking",
     "SELECT FOR UPDATE evita a venda dupla do mesmo ingresso.", badge_color=GN)

half_w = (CW - 0.30) / 2
card(sl, L, CONTENT_Y, half_w, BOTTOM - CONTENT_Y, border=RD)
T(sl, "Sem controle: o que pode acontecer", L + 0.22, CONTENT_Y + 0.18, half_w - 0.44, 0.55, sz=13.5, bold=True, color=RD, wrap=True)
bullets(sl, [
    "Worker A l\u00ea available_tickets = 1, reserva",
    "Worker B l\u00ea available_tickets = 1 ao mesmo tempo, reserva tamb\u00e9m",
], L + 0.22, CONTENT_Y + 0.78, half_w - 0.44, sz=12, dot_color=RD, text_color=TXT, dy=0.50)
R_(sl, L + 0.22, CONTENT_Y + 1.92, 0.09, 0.09, fill=RD)
T(sl, "Resultado: 2 vendas, 0 ingressos (overbooking)", L + 0.22 + 0.22, CONTENT_Y + 1.855, half_w - 0.66, 0.30, sz=12, bold=True, color=RD)

rl = L + half_w + 0.30
card(sl, rl, CONTENT_Y, half_w, BOTTOM - CONTENT_Y, border=GN)
T(sl, "Com SELECT FOR UPDATE: nossa solu\u00e7\u00e3o", rl + 0.22, CONTENT_Y + 0.18, half_w - 0.44, 0.55, sz=13.5, bold=True, color=GN, wrap=True)
code(sl, rl + 0.22, CONTENT_Y + 0.78, half_w - 0.44, 1.65,
     "-- stock-worker/src/database.py\n"
     "BEGIN;\n"
     "SELECT * FROM events\n"
     "  WHERE id = $1\n"
     "  FOR UPDATE;   -- bloqueia outros workers\n"
     "UPDATE events SET available_tickets -= $qty;\n"
     "UPDATE orders SET status = 'stock_reserved';\n"
     "COMMIT;         -- libera o lock", sz=10)
bullets(sl, [
    "Worker B tenta e bloqueia",
    "Aguarda Worker A commitar",
    "Se estoque zerou: recebe out_of_stock",
    "Sem overbooking poss\u00edvel",
], rl + 0.22, CONTENT_Y + 2.62, half_w - 0.44, sz=11.5, dot_color=GN, text_color=TXT, dy=0.40)
pagenum(sl, 21)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Demonstração Prática: TicketLab
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
head(sl, "DEMONSTRA\u00c7\u00c3O", "TicketLab: Sistema de Venda de Ingressos",
     "Materializa todos os conceitos vistos: o RabbitMQ \u00e9 o protagonista, n\u00e3o infraestrutura oculta.", badge_color=OR)

pipe_t = CONTENT_Y + 0.05
pipe_items = [("Frontend", SUBT, "React"), ("Producer", TL, "FastAPI"), ("RabbitMQ", OR, "broker"),
              ("payment", GN, "worker"), ("stock", GN, "worker"), ("notification", GN, "worker")]
box_w6 = (CW - 5 * 0.14) / 6
pipeline(sl, L, pipe_t, box_w6, 0.72, 0.14, pipe_items, title_sz=11.5)

badges_t = pipe_t + 0.72 + 0.40
status_items = [("pending", MUT), ("processing", AM), ("payment_approved", AM), ("stock_reserved", AM), ("confirmed", GN)]
box_w5 = (CW - 4 * 0.14) / 5
pipeline(sl, L, badges_t, box_w5, 0.42, 0.14, status_items, title_sz=10)

lower_t = badges_t + 0.42 + 0.32
left_w = 5.55
card(sl, L, lower_t, left_w, BOTTOM - lower_t)
T(sl, "Conceitos demonstrados", L + 0.22, lower_t + 0.18, left_w - 0.44, 0.32, sz=13.5, bold=True, color=OR)
bullets(sl, [
    "Perfect Links: filas durable + PERSISTENT + ACK manual",
    "Crash-stop: make kill-payment, sem perda de mensagens",
    "Race condition: SELECT FOR UPDATE, zero overbooking",
    "DLQ: make send-failures (30% falha) -> make inspect-dlq",
    "Escalabilidade: make scale-workers N=5",
], L + 0.22, lower_t + 0.62, left_w - 0.44, sz=11.5, dot_color=OR, text_color=TXT, dy=0.46)

rl = L + left_w + 0.30
rw = RM - rl
placeholder(sl, rl, lower_t, rw, BOTTOM - lower_t, "screenshot",
            "Frontend TicketLab: tela principal com eventos",
            "localhost:3000: capturar a Home com os cards de eventos (Coldplay, Lollapalooza etc.)")
pagenum(sl, 22)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — Demo: Saga Coreografada
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
head(sl, "DEMO", "Demo: Saga Coreografada",
     "Um pedido passa por 4 etapas; cada uma \u00e9 um worker independente.", badge_color=OR)

steps23 = [
    (1, OR, "Compra", ["Formul\u00e1rio preenchido -> POST /orders.", "API retorna status: pending de imediato.",
                        "Publica em ticket-sales (key=payment)."],
     "localhost:3000/buy: formul\u00e1rio preenchido antes de clicar"),
    (2, TL, "Pagamento", ["payment-worker consome payment-queue.", "Valida -> publica em payment-approved.",
                          "Status: processing."],
     "localhost:3000/orders: status processing"),
    (3, AM, "Estoque", ["stock-worker consome stock-queue.", "SELECT FOR UPDATE reserva o ingresso.",
                        "Status: stock_reserved."],
     "Mesma sequ\u00eancia, capturar em seguida"),
    (4, GN, "Confirma\u00e7\u00e3o", ["notification-worker gera ticket_code (UUID).", "Verifica idempot\u00eancia.",
                                  "Status: confirmed. Ingresso na m\u00e3o do comprador."],
     "localhost:3000/orders: detalhe com status confirmed"),
]
cw4 = (CW - 3 * 0.20) / 4
top_h = 1.95
for i, (num, color, title, desc, shot) in enumerate(steps23):
    cx = L + i * (cw4 + 0.20)
    stepcard(sl, cx, CONTENT_Y, cw4, top_h, num, color, title, desc, num_sz=30)
    ph_t = CONTENT_Y + top_h + 0.18
    placeholder(sl, cx, ph_t, cw4, BOTTOM - ph_t, "screenshot", shot)
pagenum(sl, 23)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 24 — Demo: Falhas na Prática
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
head(sl, "DEMO", "Demo: Falhas na Pr\u00e1tica", "Tr\u00eas cen\u00e1rios: o sistema \u00e9 resiliente.", badge_color=RD)

scenarios = [
    ("Dead Letter Queue", RD, "make send-failures\nmake inspect-dlq",
     "30% dos pedidos com simulate_failure. 3 retries, depois NACK -> DLQ.",
     "RabbitMQ UI: dead-letter-queue com N mensagens (localhost:15672)"),
    ("Crash-stop Worker", AM, "make send-1000\nmake kill-payment\nmake logs-payment",
     "Outro worker assume as mensagens unacked automaticamente.",
     "Grafana: consumer count caindo de 2 para 1 (localhost:3001)"),
    ("Crash-recovery Broker", GN, "make restart-rabbit",
     "Mensagens PERSISTENT no disco. Workers reconectam com connect_robust.",
     "RabbitMQ UI ap\u00f3s restart: filas intactas (localhost:15672)"),
]
cw3 = (CW - 2 * 0.24) / 3
top_h = 1.95
for i, (title, color, cmds, desc, shot) in enumerate(scenarios):
    cx = L + i * (cw3 + 0.24)
    card(sl, cx, CONTENT_Y, cw3, top_h)
    T(sl, title, cx + 0.20, CONTENT_Y + 0.18, cw3 - 0.40, 0.55, sz=13.5, bold=True, color=color, wrap=True)
    code(sl, cx + 0.20, CONTENT_Y + 0.72, cw3 - 0.40, 0.62, cmds, sz=9.5)
    PARA(sl, [desc], cx + 0.20, CONTENT_Y + 1.40, cw3 - 0.40, 0.55, sz=10, color=MUT, lh=1.2)
    ph_t = CONTENT_Y + top_h + 0.18
    placeholder(sl, cx, ph_t, cw3, BOTTOM - ph_t, "screenshot", shot)
pagenum(sl, 24)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 25 — Conclusão  (tema escuro, espelha a capa)
# ════════════════════════════════════════════════════════════════════════════
sl = new_slide()
bgfill(sl, CV)
badge(sl, "CONCLUS\u00c3O", color=OR)
T(sl, "Conclus\u00e3o", L - 0.02, 0.88, CW, 0.75, sz=33, bold=True, color=WH, font=SERIF, anchor='t')
T(sl, "O que aprendemos sobre o RabbitMQ em Sistemas Distribu\u00eddos", L, 1.68, CW, 0.40, sz=14, italic=True, color=CVSUB)

qa = [
    ("Onde \u00e9 usado?", "E-commerce, IoT, bancos, microservi\u00e7os: desacoplamento temporal."),
    ("Funcionalidades?", "Exchanges, durable queues, PERSISTENT, ACK, DLX, QoS, plugins."),
    ("Falhas suportadas?", "Crash-stop, crash-recovery, retry+DLQ. N\u00e3o suporta Byzantine."),
    ("Confiabilidade?", "Perfect Links (at-least-once) + aproxima\u00e7\u00e3o de Reliable Broadcast."),
    ("Ordena\u00e7\u00e3o?", "FIFO com 1 consumer. N consumers requer particionamento."),
    ("Implementa\u00e7\u00e3o?", "AMQP 0-9-1, Erlang/Mnesia, protocolo ACK, SELECT FOR UPDATE."),
]
tbl_t = 2.28
col1_w, col2_w = 3.55, CW - 3.55 - 0.18
header_h, row_h = 0.42, 0.385
R_(sl, L, tbl_t, col1_w, header_h, fill=OR)
T(sl, "QUEST\u00c3O DO PROFESSOR", L + 0.18, tbl_t, col1_w - 0.18, header_h, sz=10.5, bold=True, color=CV, anchor='ctr')
R_(sl, L + col1_w + 0.18, tbl_t, col2_w, header_h, fill=OR)
T(sl, "RESPOSTA", L + col1_w + 0.18 + 0.18, tbl_t, col2_w - 0.18, header_h, sz=10.5, bold=True, color=CV, anchor='ctr')
y = tbl_t + header_h + 0.07
ROWDARK = C("24232C")
for i, (q, a) in enumerate(qa):
    rowfill = ROWDARK if i % 2 == 0 else CV
    R_(sl, L, y, col1_w, row_h, fill=rowfill)
    T(sl, q, L + 0.18, y, col1_w - 0.30, row_h, sz=11.5, bold=True, color=OR, anchor='ctr', wrap=True)
    R_(sl, L + col1_w + 0.18, y, col2_w, row_h, fill=rowfill)
    T(sl, a, L + col1_w + 0.18 + 0.18, y, col2_w - 0.32, row_h, sz=11, color=CVSUB, anchor='ctr', wrap=True)
    y += row_h + 0.07

y += 0.10
T(sl, "Pr\u00f3ximos passos para produ\u00e7\u00e3o", L, y, CW, 0.30, sz=14, bold=True, color=OR)
nexts = [
    "Publisher confirms: Reliable Broadcast formal (saindo do BEB)",
    "Outbox Pattern: zero perda durante parti\u00e7\u00e3o producer-broker",
    "Quorum Queues (Raft): aproxima\u00e7\u00e3o de URB nativa no RabbitMQ 3.8+",
    "Kafka + partition keys: ordem total por evento",
]
bullets(sl, nexts, L, y + 0.40, CW, sz=11.5, dot_color=OR, text_color=CVSUB, dy=0.34)

T(sl, "Obrigado. Ficamos \u00e0 disposi\u00e7\u00e3o para perguntas.", L, 7.05, CW, 0.34, sz=13, bold=True,
  color=OR, align=PP_ALIGN.CENTER, anchor='ctr')


# ════════════════════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"OK: {OUTPUT}  ({len(prs.slides)} slides)")
